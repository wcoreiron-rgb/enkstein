"""Local rendering of Office documents for governed Cowork file changes.

A Brain -- browser chat session, CLI subscription, or local model -- can only
hand back text. It cannot produce a real ``.docx`` or ``.pptx`` binary. Before
this module existed, asking for "a doc and a deck" produced changes whose
``content`` was markdown written verbatim into ``report.docx``, which Word then
refused to open.

Enkstein instead treats the Brain's text as the *source* for the document and
renders the real Office binary locally in the backend. No model output is
executed; only structured text is read.
"""

from __future__ import annotations

import csv
import io
import posixpath
import re
import zipfile

# Office formats Enkstein renders locally from Brain-provided markdown.
OFFICE_EXTENSIONS = {"docx", "pptx", "xlsx"}

OFFICE_MIME_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

# Rendering is bounded so a runaway model response cannot build an enormous
# document in memory.
_MAX_SOURCE_CHARS = 400_000
_MAX_BLOCKS = 4_000
_MAX_SLIDES = 200
_MAX_ROWS = 5_000

_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET = re.compile(r"^\s*[-*+]\s+(.*)$")
_NUMBERED = re.compile(r"^\s*\d+[.)]\s+(.*)$")
_FENCE = re.compile(r"^\s*```")
_SLIDE_BREAK = re.compile(r"^\s*(---|===|\*\*\*)\s*$")
_EMPHASIS = re.compile(r"(\*\*|__|\*|_|`)")


class OfficeRenderError(RuntimeError):
    """The Brain's text could not be turned into a valid Office document."""


def office_extension(path: str) -> str | None:
    """Return the lowercased Office extension for ``path``, if it has one."""
    extension = posixpath.splitext(path)[1].lstrip(".").lower()
    return extension if extension in OFFICE_EXTENSIONS else None


def is_office_path(path: str) -> bool:
    return office_extension(path) is not None


# Text-bearing formats whose bytes can be scanned directly for secrets. OOXML
# and archives are ZIP containers, so their text lives in compressed members.
_ZIP_CONTAINER_EXTENSIONS = {"docx", "pptx", "xlsx", "zip"}
_MAX_SCAN_BYTES = 2_000_000
_MAX_ZIP_MEMBERS = 200


def extract_scannable_text(name: str, payload: bytes) -> str:
    """Best-effort readable text from a provider-generated file, for DLP scanning.

    A harvested download is opaque bytes, so a naive scan would miss an API key
    sitting inside a generated spreadsheet or script. Office/ZIP payloads are
    unpacked in memory and their text members concatenated; anything else is
    decoded as UTF-8. Returns "" when nothing readable can be recovered, which
    the caller must treat as "unscannable" rather than "clean".
    """
    extension = posixpath.splitext(name)[1].lstrip(".").lower()
    if extension in _ZIP_CONTAINER_EXTENSIONS:
        collected: list[str] = []
        total = 0
        try:
            with zipfile.ZipFile(io.BytesIO(payload)) as archive:
                for info in archive.infolist()[:_MAX_ZIP_MEMBERS]:
                    if info.is_dir() or info.file_size > _MAX_SCAN_BYTES:
                        continue
                    # Guard against a zip bomb: stop once enough text is read
                    # to scan meaningfully.
                    if total >= _MAX_SCAN_BYTES:
                        break
                    try:
                        with archive.open(info) as member:
                            raw = member.read(_MAX_SCAN_BYTES - total)
                    except Exception:
                        continue
                    total += len(raw)
                    collected.append(raw.decode("utf-8", errors="ignore"))
        except (zipfile.BadZipFile, OSError):
            return ""
        return "\n".join(collected)
    return payload[:_MAX_SCAN_BYTES].decode("utf-8", errors="ignore")


def _clean(text: str) -> str:
    """Strip inline markdown emphasis so it does not leak into Office text."""
    return _EMPHASIS.sub("", text).strip()


def _blocks(source: str) -> list[tuple[str, int, str]]:
    """Parse markdown-ish text into ``(kind, level, text)`` blocks.

    ``kind`` is one of ``heading``, ``bullet``, ``break``, or ``text``. Fenced
    code blocks are preserved verbatim as ``text`` so scripts embedded in a
    document survive rendering.
    """
    blocks: list[tuple[str, int, str]] = []
    in_fence = False
    for line in source.splitlines():
        if len(blocks) >= _MAX_BLOCKS:
            break
        if _FENCE.match(line):
            in_fence = not in_fence
            continue
        if in_fence:
            blocks.append(("text", 0, line.rstrip()))
            continue
        stripped = line.strip()
        if not stripped:
            continue
        if _SLIDE_BREAK.match(stripped):
            blocks.append(("break", 0, ""))
            continue
        heading = _HEADING.match(stripped)
        if heading:
            text = _clean(heading.group(2))
            if text:
                blocks.append(("heading", len(heading.group(1)), text))
            continue
        bullet = _BULLET.match(line) or _NUMBERED.match(line)
        if bullet:
            text = _clean(bullet.group(1))
            if text:
                blocks.append(("bullet", 0, text))
            continue
        blocks.append(("text", 0, _clean(stripped)))
    return blocks


def _render_docx(source: str) -> bytes:
    from docx import Document

    document = Document()
    for kind, level, text in _blocks(source):
        if kind == "heading":
            document.add_heading(text, level=min(level, 9))
        elif kind == "bullet":
            document.add_paragraph(text, style="List Bullet")
        elif kind == "break":
            document.add_page_break()
        elif text:
            document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _render_pptx(source: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    title_and_content = presentation.slide_layouts[1]
    state: dict[str, object] = {"body": None, "count": 0, "first_used": True}

    def start_slide(title: str) -> None:
        if int(state["count"]) >= _MAX_SLIDES:
            state["body"] = None
            return
        state["count"] = int(state["count"]) + 1
        slide = presentation.slides.add_slide(title_and_content)
        slide.shapes.title.text = title[:180]
        frame = slide.placeholders[1].text_frame
        frame.clear()
        state["body"] = frame
        # A cleared frame keeps one empty paragraph; reuse it for the first
        # line rather than leaving a blank bullet at the top of the slide.
        state["first_used"] = False

    def add_line(text: str) -> None:
        frame = state["body"]
        if frame is None or not text:
            return
        if not state["first_used"]:
            paragraph = frame.paragraphs[0]
            state["first_used"] = True
        else:
            paragraph = frame.add_paragraph()
        paragraph.text = text[:500]

    for kind, level, text in _blocks(source):
        if kind == "break":
            state["body"] = None
        elif kind == "heading" and level <= 2:
            start_slide(text)
        elif kind == "heading":
            if state["body"] is None:
                start_slide(text)
            else:
                add_line(text)
        else:
            if state["body"] is None:
                start_slide("Overview")
            add_line(text)

    if int(state["count"]) == 0:
        blank = presentation.slides.add_slide(presentation.slide_layouts[5])
        blank.shapes.title.text = "Untitled"
        textbox = blank.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        textbox.text_frame.text = source.strip()[:2000]

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _spreadsheet_rows(source: str) -> list[list[str]]:
    """Read CSV or markdown-table text into rows."""
    lines = [line for line in source.splitlines() if line.strip()]
    pipe_rows = [line for line in lines if line.strip().startswith("|")]
    if len(pipe_rows) >= 2:
        rows: list[list[str]] = []
        for line in pipe_rows[:_MAX_ROWS]:
            cells = [_clean(cell) for cell in line.strip().strip("|").split("|")]
            # Markdown separator rows (---|:--:) carry no data.
            if cells and all(set(cell) <= set("-: ") for cell in cells if cell):
                continue
            rows.append(cells)
        if rows:
            return rows
    reader = csv.reader(io.StringIO("\n".join(lines[:_MAX_ROWS])))
    return [row for row in reader if row]


def _render_xlsx(source: str) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    for row in _spreadsheet_rows(source):
        sheet.append([cell[:5000] for cell in row])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def render_office_document(path: str, source: str) -> bytes:
    """Render Brain-provided text into a real Office binary for ``path``.

    Raises :class:`OfficeRenderError` when the extension is unsupported, the
    source is oversized, or the rendering library is unavailable, so the caller
    can skip the change rather than write an unopenable file.
    """
    extension = office_extension(path)
    if extension is None:
        raise OfficeRenderError("Unsupported Office document type")
    if len(source) > _MAX_SOURCE_CHARS:
        raise OfficeRenderError("Document source exceeds the supported size")
    renderers = {"docx": _render_docx, "pptx": _render_pptx, "xlsx": _render_xlsx}
    try:
        return renderers[extension](source)
    except OfficeRenderError:
        raise
    except ImportError as exc:
        raise OfficeRenderError("Office rendering support is not installed") from exc
    except Exception as exc:  # noqa: BLE001 - any library failure is a skip, not a crash
        raise OfficeRenderError("The document could not be rendered") from exc
